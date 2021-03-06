def get_FilePaths(directory, extension=['']):
    """
    /*-----------------------------------------------------*/
    description:
        This function will generate the file names in a directory 
        tree by walking the tree either top-down or bottom-up. For each 
        directory in the tree rooted at directory top (including top itself), 
        it yields a 3-tuple (dirpath, dirnames, filenames).
    /*-----------------------------------------------------*/
    args: 
        directory: current path that is to be searched
        extension(s): specific file extension to be found.
                    defaults to empty, which will give all files
                    found in tree
    returns:
        file_paths: list of strings with each absolute path
    /*----------------------------------------------------*/
    """
    import os 
    file_paths = [] # init of list of strings to be returned

    # walk the tree
    for root, directories, files in os.walk(directory):
        for filename in files:
            # join the two strings in order to form the full filepath 
            filepath = os.path.join(root, filename)
            # check the filepath's extension, if wanted then add to list
            # if arg:extension is blank, then all files found in tree
            # will be added
            for ext in extension:
                if filepath.lower().endswith(ext) and 'analyzed' not in filepath:
                    file_paths.append(filepath) # add it to the list.
    return file_paths # self-explanatory

def generate_PicCoords(x, y):
    """
    /*-----------------------------------------------------*/
    description:
        Generate an array for placing images on a ppt slide
        in a 2 by 2 manner.
    /*-----------------------------------------------------*/
    args: 
        x: the slide width
        y: the slide height
    returns:
        multi-dim array of int coordinates
        example:
            [ [0, 0], [4572000, 0], 
            [0, 3429000], [4572000, 3429000] ]
        accessing elements:
            [0][0] = 0; [0][1] = 0;
            [1][0] = 4572000; [1][1] = 0;
            [2][0] = 0; [2][1] = 3429000;
            [3][0] = 4572000; 3429000;
        ** this is just an example and coordinates will
        vary depending on default slide width, height, 
        and the image.
    /*----------------------------------------------------*/
    """
    return [ [0,0], [x, 0], [0, y], [x, y] ]


def generate_Groups(iterable, n, fillvalue=None):
    """
    /*-----------------------------------------------------*/
    description:
        Take an iterable object and create groups of length 
        n.
    /*-----------------------------------------------------*/
    args: 
        iterable: an object to be iterated over, can be list,
            string, tuple, etc...

        n: the length of each group

        fillvalue: (defaults to None) this is used when
            len(iterable) % n != 0, then fillvalue completes
            the remaining spaces in the group (see examples 
            below)
    returns:
        an iterator object

        ** Example:
        generate_Groups('ABCDEFG', 3, 'x') produces
        --> ('A', 'B', 'C'), ('D', 'E', 'F'), ('G', 'x', 'x')

        ** Example:
        generate_Groups([1,2,3,4,5], 2) produces
        --> (1, 2), (3, 4), (5, None)
    /*----------------------------------------------------*/
    """
    from itertools import izip_longest
    args = [iter(iterable)] * n
    return izip_longest(fillvalue=fillvalue, *args)

def generate_SortedFiles(f_list):
    """
    /*-----------------------------------------------------*/
    description:
        Sort a list of files by date created. This defaults 
        to sorting by [newest, ..., oldest] (this is 
        quasi-FIFO), but to change this to the reverse
        (this is quasi-LIFO) then add: f_list.reverse() at 
        the end of the function
    /*-----------------------------------------------------*/
    args: 
        f_list: a list of files
    returns:
        No return. f_list is passed by reference.
    /*----------------------------------------------------*/
    """
    import os
    f_list.sort(key=lambda x: os.path.getmtime(x))


if __name__=="__main__":
    import pptx
    import pptx.util
    import scipy.misc
    import os
    import time 

    # image file extensions
    img_extensions = ['.png', '.jpg', '.jpeg', '.tiff']
    # absolute path
    cwd = os.path.dirname(os.path.abspath(__file__)) 

    # init new presentation
    prs = pptx.Presentation()
    # set slide height @ 4:3
    prs.slide_height = 6858000

    # put image paths into list
    img_file_path = get_FilePaths(cwd, img_extensions)
    # sort the images [oldest, ..., newest]

    # group images in groups of n and loop (n = 4)
    group_imgs = generate_Groups(img_file_path, 4)
    for group in group_imgs:
        # add a blank slide for each group of n
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # loop over image in each group of n
        for i in range( len(group) ):
            # img as 2d numpy array
            # if group has nonetype in it then
            # attribute error is raised and we skip 
            # that image
            try:
                img = scipy.misc.imread( group[i] )
            except AttributeError:
                continue # go to top of loop
            # set width
            width = int(prs.slide_width * 0.5)
            # set height based on img dims
            height = int(width * img.shape[0]/img.shape[1])
            # generate coordinate array
            xy = generate_PicCoords(width, height)
            #xy[i][0] = width coordinate; xy[i][1] = height coordinate
            # add img to slide (img is accessed as group[i])
            pic = slide.shapes.add_picture(group[i], xy[i][0], xy[i][1], height)

    # save ppt presentation
    try:
        prs.save('test.pptx')
    except IOError:
        print
        print "Close PowerPoint and re-run"
        print
